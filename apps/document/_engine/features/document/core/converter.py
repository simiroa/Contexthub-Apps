import os
import sys
import threading
from pathlib import Path
import logging
import json
import time
import zipfile
import shutil

# Optional Dependencies
try: import pymupdf; HAS_PYMUPDF = True
except ImportError: HAS_PYMUPDF = False

try: from pdf2docx import Converter as PDF2DOCXConverter; HAS_PDF2DOCX = True
except ImportError: HAS_PDF2DOCX = False

try: import comtypes.client; HAS_COMTYPES = True
except ImportError: HAS_COMTYPES = False

try: import markdown; HAS_MARKDOWN = True
except ImportError: HAS_MARKDOWN = False

try: from xhtml2pdf import pisa; HAS_PISA = True
except ImportError: HAS_PISA = False

try: from PIL import Image; HAS_PIL = True
except ImportError: HAS_PIL = False

logger = logging.getLogger("doc_converter_backend")

class DocumentConverter:
    """
    Core backend for document conversion.
    Supports:
    - PDF <-> Word/Excel/PPTX (Office required for Office -> PDF)
    - PDF <-> Image (JPG/PNG)
    - Markdown -> PDF/HTML
    - PDF -> EPUB
    """
    
    def __init__(self):
        self._is_cancelled = False

    def cancel(self):
        self._is_cancelled = True

    def convert(self, source_path, target_ext, output_dir=None, options=None):
        """
        Main entry point for conversion.
        :param source_path: Path to the source file.
        :param target_ext: Target extension (e.g., '.pdf', '.docx', '.png').
        :param output_dir: Directory to save the output. Defaults to source dir.
        :param options: Dict of additional options (DPI, separate_pages, etc.)
        """
        self._is_cancelled = False
        source = Path(source_path)
        if not source.exists():
            raise FileNotFoundError(f"Source file not found: {source}")
            
        target_ext = target_ext.lower()
        if not target_ext.startswith('.'):
            target_ext = f".{target_ext}"
            
        if not output_dir:
            output_dir = source.parent
        else:
            output_dir = Path(output_dir)
            output_dir.mkdir(parents=True, exist_ok=True)
            
        source_ext = source.suffix.lower()
        
        # Dispatch based on target/source
        if target_ext == '.images':
            if source_ext == '.pdf':
                return self._extract_images_from_pdf(source, output_dir)
            elif source_ext in ['.pptx', '.docx', '.xlsx']:
                return self._extract_office_images(source, output_dir)
            else:
                raise ValueError(f"Image extraction not supported for {source_ext}")
        
        if source_ext == '.pdf':
            return self._convert_from_pdf(source, target_ext, output_dir, options or {})
        elif target_ext == '.pdf':
            return self._convert_to_pdf(source, output_dir, options or {})
        elif source_ext == '.md':
            return self._convert_markdown(source, target_ext, output_dir, options or {})
        else:
            raise ValueError(f"Unsupported conversion: {source_ext} to {target_ext}")

    def _convert_from_pdf(self, pdf_path, target_ext, output_dir, options):
        """PDF -> Other formats"""
        if not HAS_PYMUPDF:
            raise ImportError("PyMuPDF (pymupdf) is required for PDF operations.")
            
        base_name = pdf_path.stem
        
        if target_ext in ['.docx', '.doc']:
            if not HAS_PDF2DOCX:
                raise ImportError("pdf2docx is required for PDF to Word conversion.")
            out_path = output_dir / f"{base_name}.docx"
            cv = PDF2DOCXConverter(str(pdf_path))
            try:
                cv.convert(str(out_path))
            finally:
                cv.close()
            return [out_path]
            
        elif target_ext in ['.xlsx', '.xls', '.csv']:
            # For simplicity, extract tables to CSV/XLSX if pandas is available, 
            # or just extract text as a fallback.
            return self._extract_tables(pdf_path, output_dir, format=target_ext[1:])
            
        elif target_ext == '.pptx':
            # PDF to PPTX is complex. We'll do a simple page-as-image-in-slide approach or similar.
            # For now, let's mark it as limited or use an external tool if available.
            return self._pdf_to_pptx_basic(pdf_path, output_dir)
            
        elif target_ext in ['.png', '.jpg', '.jpeg']:
            dpi = options.get('dpi', 300)
            separate_pages = options.get('separate_pages', True)
            return self._pdf_to_images(pdf_path, target_ext[1:], output_dir, dpi)
            
        elif target_ext == '.epub':
            return self._pdf_to_epub(pdf_path, output_dir)
            
        elif target_ext == '.html':
            out_path = output_dir / f"{base_name}.html"
            self._pdf_to_html(pdf_path, out_path)
            return [out_path]
            
        elif target_ext == '.md':
            out_path = output_dir / f"{base_name}.md"
            self._pdf_to_md(pdf_path, out_path)
            return [out_path]
            
        else:
            raise ValueError(f"Unsupported target format from PDF: {target_ext}")

    def _convert_to_pdf(self, source_path, output_dir, options):
        """Other formats -> PDF"""
        source_ext = source_path.suffix.lower()
        out_path = output_dir / f"{source_path.stem}.pdf"
        
        if source_ext in ['.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt']:
            return self._office_to_pdf(source_path, out_path)
            
        elif source_ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
            return self._images_to_pdf([source_path], out_path)
            
        elif source_ext == '.md':
            return self._markdown_to_pdf(source_path, out_path)
            
        else:
            raise ValueError(f"Unsupported source format for PDF conversion: {source_ext}")

    def _convert_markdown(self, md_path, target_ext, output_dir, options):
        """Markdown -> HTML/PDF"""
        if not HAS_MARKDOWN:
            raise ImportError("markdown library is required for Markdown conversions.")
            
        base_name = md_path.stem
        
        if target_ext == '.html':
            out_path = output_dir / f"{base_name}.html"
            with open(md_path, 'r', encoding='utf-8') as f:
                content = f.read()
            html = markdown.markdown(content, extensions=['extra', 'tables', 'toc'])
            # Wrap in basic HTML5 boilerplate
            full_html = f"<!DOCTYPE html><html><head><meta charset='utf-8'></head><body>{html}</body></html>"
            with open(out_path, 'w', encoding='utf-8') as f:
                f.write(full_html)
            return [out_path]
            
        elif target_ext == '.pdf':
            out_path = output_dir / f"{base_name}.pdf"
            return self._markdown_to_pdf(md_path, out_path)
            
        else:
            raise ValueError(f"Unsupported target format from Markdown: {target_ext}")

    # --- Worker Methods ---

    def _office_to_pdf(self, source_path, out_path):
        """Convert Office docs to PDF using comtypes (Windows only)"""
        if not HAS_COMTYPES:
            raise ImportError("comtypes is required for Office to PDF conversion on Windows.")
            
        source_ext = source_path.suffix.lower()
        abs_source = str(source_path.absolute())
        abs_out = str(out_path.absolute())
        
        if source_ext in ['.docx', '.doc']:
            app = comtypes.client.CreateObject("Word.Application")
            app.Visible = False
            try:
                doc = app.Documents.Open(abs_source)
                doc.ExportAsFixedFormat(abs_out, 17) # 17 = wdExportFormatPDF
                doc.Close()
            finally:
                app.Quit()
                
        elif source_ext in ['.xlsx', '.xls']:
            app = comtypes.client.CreateObject("Excel.Application")
            app.Visible = False
            try:
                wb = app.Workbooks.Open(abs_source)
                wb.ExportAsFixedFormat(0, abs_out) # 0 = xlTypePDF
                wb.Close()
            finally:
                app.Quit()
                
        elif source_ext in ['.pptx', '.ppt']:
            app = comtypes.client.CreateObject("PowerPoint.Application")
            # PowerPoint visible=False is slightly different
            try:
                pres = app.Presentations.Open(abs_source, WithWindow=False)
                pres.SaveAs(abs_out, 32) # 32 = ppSaveAsPDF
                pres.Close()
            finally:
                app.Quit()
                
        return [out_path]

    def _pdf_to_images(self, pdf_path, fmt, output_dir, dpi):
        """Convert PDF pages to images"""
        doc = pymupdf.open(str(pdf_path))
        zoom = dpi / 72.0
        mat = pymupdf.Matrix(zoom, zoom)
        
        out_paths = []
        for i, page in enumerate(doc):
            if self._is_cancelled: break
            pix = page.get_pixmap(matrix=mat, alpha=False)
            out_path = output_dir / f"{pdf_path.stem}_p{i+1:03d}.{fmt}"
            pix.save(str(out_path))
            out_paths.append(out_path)
            
        doc.close()
        return out_paths

    def _images_to_pdf(self, image_paths, out_path):
        """Convert one or more images to a single PDF"""
        if not HAS_PIL:
            raise ImportError("Pillow (PIL) is required for Image to PDF conversion.")
            
        imgs = []
        for p in image_paths:
            img = Image.open(p)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            imgs.append(img)
            
        if not imgs:
            return []
            
        imgs[0].save(out_path, save_all=True, append_images=imgs[1:])
        return [out_path]

    def _markdown_to_pdf(self, md_path, out_path):
        """Convert Markdown to PDF via HTML"""
        if not HAS_MARKDOWN or not HAS_PISA:
            raise ImportError("markdown and xhtml2pdf are required for MD to PDF conversion.")
            
        with open(md_path, 'r', encoding='utf-8') as f:
            md_text = f.read()
            
        html_content = markdown.markdown(md_text, extensions=['extra', 'tables', 'toc'])
        
        # Add basic styling for PDF
        styled_html = f"""
        <html>
        <head>
            <meta charset="utf-8">
            <style>
                body {{ font-family: Helvetica, Arial, sans-serif; line-height: 1.6; padding: 20px; }}
                h1, h2, h3 {{ color: #333; }}
                table {{ border-collapse: collapse; width: 100%; border: 1px solid #ccc; }}
                th, td {{ border: 1px solid #ccc; padding: 8px; text-align: left; }}
                th {{ background-color: #f2f2f2; }}
            </style>
        </head>
        <body>{html_content}</body>
        </html>
        """
        
        with open(out_path, "wb") as f:
            pisa_status = pisa.CreatePDF(styled_html, dest=f)
            
        if pisa_status.err:
            raise RuntimeError(f"Failed to generate PDF from Markdown: {pisa_status.err}")
            
        return [out_path]

    def _extract_tables(self, pdf_path, output_dir, format='csv'):
        """Extract tables from PDF"""
        doc = pymupdf.open(str(pdf_path))
        out_paths = []
        
        for i, page in enumerate(doc):
            if self._is_cancelled: break
            tabs = page.find_tables()
            for j, tab in enumerate(tabs):
                out_path = None
                if format == 'csv':
                    out_path = output_dir / f"{pdf_path.stem}_p{i+1}_t{j+1}.csv"
                    import csv
                    with open(out_path, 'w', newline='', encoding='utf-8-sig') as f:
                        writer = csv.writer(f)
                        writer.writerows(tab.extract())
                elif format == 'xlsx':
                    out_path = output_dir / f"{pdf_path.stem}_p{i+1}_t{j+1}.xlsx"
                    try:
                        import openpyxl
                        wb = openpyxl.Workbook()
                        ws = wb.active
                        for row_data in tab.extract():
                            ws.append(row_data)
                        wb.save(str(out_path))
                    except ImportError:
                        # Fallback to CSV
                        out_path = out_path.with_suffix('.csv')
                        import csv
                        with open(out_path, 'w', newline='', encoding='utf-8-sig') as f:
                            writer = csv.writer(f)
                            writer.writerows(tab.extract())
                
                if out_path:
                    out_paths.append(out_path)
                
        doc.close()
        return out_paths

    def _pdf_to_epub(self, pdf_path, output_dir):
        """Basic PDF to EPUB conversion (Text based)"""
        # This is very basic text extraction. Real EPUB conversion is complex.
        doc = pymupdf.open(str(pdf_path))
        out_path = output_dir / f"{pdf_path.stem}.epub"
        
        # For now, let's just use pymupdf's text extraction as a placeholder
        # In a real app, you might use ebooklib.
        text = ""
        for page in doc:
            text += f"<h1>Page {page.number + 1}</h1>\n"
            text += page.get_text("html")
            
        # Very simple "EPUB" structure (actually just a packaged HTML if we were rigorous)
        # For this tool, we'll extract as HTML and note it's for EPUB.
        # REAL implementation would need ebooklib.
        try:
            from ebooklib import epub
            book = epub.EpubBook()
            book.set_title(pdf_path.stem)
            book.set_language('en')
            
            chapters = []
            for i, page in enumerate(doc):
                c = epub.EpubHtml(title=f'Page {i+1}', file_name=f'page_{i+1}.xhtml', lang='en')
                c.content = page.get_text("html")
                book.add_item(c)
                chapters.append(c)
                
            book.toc = tuple(chapters)
            book.add_item(epub.EpubNcx())
            book.add_item(epub.EpubNav())
            book.spine = ['nav'] + chapters
            
            epub.write_epub(str(out_path), book, {})
        except ImportError:
            # Fallback: Just save as HTML but rename to epub? No, that's bad.
            # Warn the user or just save as text/html.
            raise ImportError("ebooklib is required for PDF to EPUB conversion.")
            
        doc.close()
        return [out_path]

    def _pdf_to_html(self, pdf_path, out_path):
        doc = pymupdf.open(str(pdf_path))
        html = ""
        for page in doc:
            html += page.get_text("html")
        with open(out_path, 'w', encoding='utf-8') as f:
            f.write(html)
        doc.close()

    def _pdf_to_md(self, pdf_path, out_path):
        # Prefer pymupdf4llm if available
        try:
            import pymupdf4llm
            md_text = pymupdf4llm.to_markdown(str(pdf_path))
        except ImportError:
            doc = pymupdf.open(str(pdf_path))
            md_text = ""
            for page in doc:
                md_text += f"## Page {page.number + 1}\n\n"
                md_text += page.get_text() + "\n\n"
            doc.close()
            
        with open(out_path, 'w', encoding='utf-8') as f:
            f.write(md_text)

    def _pdf_to_pptx_basic(self, pdf_path, output_dir):
        """Placeholder for PDF to PPTX"""
        # Usually done by converting each page to image and placing on slide
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            raise ImportError("python-pptx is required for PDF to PPTX conversion.")
            
        doc = pymupdf.open(str(pdf_path))
        prs = Presentation()
        
        # Set slide dimensions to match PDF aspect ratio if possible, or just default
        # For simplicity, use images
        for i, page in enumerate(doc):
            if self._is_cancelled: break
            pix = page.get_pixmap(dpi=150)
            img_path = output_dir / f"temp_page_{i}.png"
            pix.save(str(img_path))
            
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
            slide.shapes.add_picture(str(img_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            os.remove(img_path)
            
        out_path = output_dir / f"{pdf_path.stem}.pptx"
        prs.save(str(out_path))
        doc.close()
        return [out_path]

    def _extract_images_from_pdf(self, pdf_path, output_dir):
        """Extract embedded images from PDF without rendering pages"""
        doc = pymupdf.open(str(pdf_path))
        img_dir = output_dir / f"Images_{pdf_path.stem}"
        img_dir.mkdir(exist_ok=True)
        
        extracted_paths = []
        processed_xrefs = set()
        
        for pno in range(len(doc)):
            if self._is_cancelled: break
            img_list = doc.get_page_images(pno)
            for img in img_list:
                xref = img[0]
                if xref in processed_xrefs: continue
                
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                out_path = img_dir / f"img_{xref}.{image_ext}"
                with open(out_path, "wb") as f:
                    f.write(image_bytes)
                
                extracted_paths.append(out_path)
                processed_xrefs.add(xref)
                
        doc.close()
        return extracted_paths

    def _extract_office_images(self, office_path, output_dir):
        """Extract images from PPTX/DOCX/XLSX (Zip based extraction)"""
        img_dir = output_dir / f"Images_{office_path.stem}"
        img_dir.mkdir(exist_ok=True)
        
        extracted_paths = []
        with zipfile.ZipFile(office_path, 'r') as z:
            # Office documents store images in 'word/media/', 'ppt/media/', or 'xl/media/'
            media_files = [f for f in z.namelist() if '/media/' in f]
            
            for media_file in media_files:
                if self._is_cancelled: break
                filename = os.path.basename(media_file)
                out_path = img_dir / filename
                
                with z.open(media_file) as source, open(out_path, "wb") as target:
                    shutil.copyfileobj(source, target)
                
                extracted_paths.append(out_path)
                
        return extracted_paths

if __name__ == "__main__":
    # Test
    conv = DocumentConverter()
    print("Backend ready.")
